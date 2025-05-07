import os
import glob
import comtypes.client
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pypdf import PdfWriter
from pdf2image import convert_from_path
from pptx import Presentation
from io import BytesIO
import sys

PP_FORMAT_PDF = 32

def convert_pptx_to_pdf(input_path, output_path):
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    comtypes.CoInitialize()
    powerpoint = comtypes.client.CreateObject('PowerPoint.Application')
    presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
    presentation.SaveAs(output_path, PP_FORMAT_PDF)
    presentation.Close()
    powerpoint.Quit()
    comtypes.CoUninitialize()

def create_black_page(output_path):
    c = canvas.Canvas(output_path, pagesize=letter)
    c.setFillColorRGB(0, 0, 0)
    c.rect(0, 0, letter[0], letter[1], fill=1)
    c.showPage()
    c.save()

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python PPTcompiler.py <input_folder>")
        sys.exit(1)
    input_folder = sys.argv[1]

    # List all .pptx files
    pptx_files = glob.glob(os.path.join(input_folder, "*.pptx"))
    if not pptx_files:
        print("No .pptx files found in the specified folder.")
        sys.exit(1)

    # Convert each .pptx to PDF
    for pptx_file in pptx_files:
        pdf_file = os.path.splitext(pptx_file)[0] + ".pdf"
        convert_pptx_to_pdf(pptx_file, pdf_file)

    # Create blank black PDF page
    black_page_path = os.path.join(input_folder, "black_page.pdf")
    create_black_page(black_page_path)

    # Get list of generated PDFs
    pdf_files = [os.path.splitext(os.path.basename(pptx))[0] + ".pdf" for pptx in pptx_files]
    pdf_files = [os.path.join(input_folder, pdf) for pdf in pdf_files]

    # Merge PDFs with black page in between
    merger = PdfWriter()
    for i, pdf in enumerate(pdf_files):
        merger.append(pdf)
        if i < len(pdf_files) - 1:
            merger.append(black_page_path)
    merged_pdf_path = os.path.join(input_folder, "merged.pdf")
    merger.write(merged_pdf_path)
    merger.close()

    # Convert merged PDF to images
    images = convert_from_path(merged_pdf_path)

    # Create new .pptx with each image as a slide
    prs = Presentation()
    for image in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
        img_stream = BytesIO()
        image.save(img_stream, format='PNG')
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
    final_pptx_path = os.path.join(input_folder, "final_presentation.pptx")
    prs.save(final_pptx_path)

    print(f"Final presentation saved as {final_pptx_path}")